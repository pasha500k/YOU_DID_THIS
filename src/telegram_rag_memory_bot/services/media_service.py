"""
????: media_service.py
???????? ? ??????? ? ?????: ????????? ????????, ????????? ?????,
????? ? ?????, ? ????? ????????? ?????????? ???????? ????????????.
"""

from __future__ import annotations

import csv
import json
import logging
import mimetypes
import re
import shutil
import subprocess
from pathlib import Path
from time import monotonic
from uuid import uuid4

import cv2
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from telegram_rag_memory_bot.config import Settings

LOGGER = logging.getLogger(__name__)


class MediaService:
    IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif", ".tiff", ".tif", ".heic"}
    VIDEO_EXTENSIONS = {".mp4", ".m4v", ".mov", ".webm", ".mkv", ".avi", ".mpeg", ".mpg"}
    AUDIO_EXTENSIONS = {".mp3", ".wav", ".ogg", ".oga", ".m4a", ".aac", ".flac", ".opus", ".wma"}
    TEXT_EXTENSIONS = {
        ".txt",
        ".md",
        ".markdown",
        ".py",
        ".js",
        ".ts",
        ".tsx",
        ".jsx",
        ".java",
        ".go",
        ".rs",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
        ".css",
        ".scss",
        ".html",
        ".htm",
        ".xml",
        ".json",
        ".yaml",
        ".yml",
        ".csv",
        ".log",
        ".ini",
        ".cfg",
        ".sql",
        ".srt",
        ".rtf",
    }

    def __init__(self, settings: Settings) -> None:
        self.settings = settings

    def create_work_dir(self) -> Path:
        work_dir = self.settings.media_cache_dir / uuid4().hex
        work_dir.mkdir(parents=True, exist_ok=True)
        return work_dir

    def cleanup_work_dir(self, work_dir: Path) -> None:
        shutil.rmtree(work_dir, ignore_errors=True)

    @staticmethod
    def detect_message_type(message: object) -> str:
        if getattr(message, "photo", None):
            return "image"
        if getattr(message, "video", None):
            return "video"
        if getattr(message, "audio", None) or getattr(message, "voice", None):
            return "audio"
        if getattr(message, "document", None):
            return "document"
        if getattr(message, "raw_text", None):
            return "text"
        return "unknown"

    @classmethod
    def detect_path_type(cls, file_path: Path | None) -> str:
        if file_path is None:
            return "unknown"
        suffix = file_path.suffix.lower()
        if suffix in cls.IMAGE_EXTENSIONS:
            return "image"
        if suffix in cls.VIDEO_EXTENSIONS:
            return "video"
        if suffix in cls.AUDIO_EXTENSIONS:
            return "audio"
        if suffix in cls.TEXT_EXTENSIONS or suffix in {".pdf", ".docx", ".pptx", ".xlsx", ".json", ".csv", ".html", ".htm", ".xml"}:
            return "document"
        return "document"

    @staticmethod
    def is_video_message(message: object) -> bool:
        if getattr(message, "video", None) or getattr(message, "video_note", None):
            return True
        file_obj = getattr(message, "file", None)
        mime_type = getattr(file_obj, "mime_type", None) if file_obj else None
        return bool(mime_type and mime_type.startswith("video/"))

    @staticmethod
    def get_caption(message: object) -> str:
        return (getattr(message, "message", None) or getattr(message, "raw_text", None) or "").strip()

    @staticmethod
    def get_file_name(message: object, message_type: str) -> str | None:
        file_obj = getattr(message, "file", None)
        if not file_obj:
            return None
        if getattr(file_obj, "name", None):
            return file_obj.name
        ext = getattr(file_obj, "ext", None) or ""
        return f"{message_type}{ext}" if ext else message_type

    @staticmethod
    def get_mime_type(message: object) -> str | None:
        file_obj = getattr(message, "file", None)
        return getattr(file_obj, "mime_type", None) if file_obj else None

    @staticmethod
    def get_file_size(message: object) -> int | None:
        file_obj = getattr(message, "file", None)
        return getattr(file_obj, "size", None) if file_obj else None

    @staticmethod
    def get_file_name_from_path(file_path: Path | None) -> str | None:
        return file_path.name if file_path is not None else None

    @staticmethod
    def get_mime_type_from_path(file_path: Path | None) -> str | None:
        if file_path is None:
            return None
        mime_type, _encoding = mimetypes.guess_type(str(file_path))
        return mime_type

    @staticmethod
    def get_file_size_from_path(file_path: Path | None) -> int | None:
        if file_path is None or not file_path.exists():
            return None
        return file_path.stat().st_size

    async def download_message_media(self, client: object, message: object, work_dir: Path) -> Path | None:
        if not getattr(message, "media", None):
            return None
        file_name = self.get_file_name(message, self.detect_message_type(message))
        progress_callback = self.build_transfer_progress_callback(
            action="Скачиваю файл из Telegram",
            message=message,
            target_name=file_name,
        )
        result = await client.download_media(
            message,
            file=str(work_dir),
            progress_callback=progress_callback,
        )
        if not result:
            return None
        return Path(result)

    async def download_video_message(self, client: object, message: object) -> Path | None:
        if not self.is_video_message(message):
            return None

        target_path = self._build_video_download_path(message)
        if target_path.exists():
            return target_path

        progress_callback = self.build_transfer_progress_callback(
            action="Скачиваю файл из Telegram",
            message=message,
            target_name=target_path.name,
        )
        result = await client.download_media(
            message,
            file=str(target_path),
            progress_callback=progress_callback,
        )
        if not result:
            return None
        downloaded_path = Path(result)
        return downloaded_path if downloaded_path.exists() else target_path

    def build_transfer_progress_callback(
        self,
        *,
        action: str,
        message: object,
        target_name: str | None = None,
    ) -> object:
        file_name = (target_name or self.get_file_name(message, self.detect_message_type(message)) or "file").strip()
        start_at = monotonic()
        last_percent = -1.0
        last_bytes = 0
        last_log_at = 0.0

        def callback(current: int, total: int) -> None:
            nonlocal last_percent, last_bytes, last_log_at
            if total <= 0:
                return
            percent = (current / total) * 100
            now = monotonic()
            should_log = (
                current >= total
                or last_percent < 0
                or percent - last_percent >= 1.0
                or current - last_bytes >= 5 * 1024 * 1024
                or now - last_log_at >= 5.0
            )
            if not should_log:
                return

            elapsed = max(now - start_at, 0.001)
            speed_bytes = current / elapsed
            eta_seconds = None
            if speed_bytes > 0 and current < total:
                eta_seconds = max((total - current) / speed_bytes, 0.0)

            last_percent = percent
            last_bytes = current
            last_log_at = now
            eta_text = ""
            if eta_seconds is not None:
                eta_text = f" | осталось ~{self._format_duration(eta_seconds)}"
            LOGGER.info(
                "%s: %s | %.1f%% | %s из %s | скорость %s/с%s",
                action,
                file_name,
                percent,
                self._format_bytes(current),
                self._format_bytes(total),
                self._format_bytes(int(speed_bytes)),
                eta_text,
            )

        return callback

    @staticmethod
    def _format_bytes(size_bytes: int) -> str:
        size = float(max(size_bytes, 0))
        units = ["B", "KB", "MB", "GB", "TB"]
        unit_index = 0
        while size >= 1024 and unit_index < len(units) - 1:
            size /= 1024
            unit_index += 1
        if unit_index == 0:
            return f"{int(size)} {units[unit_index]}"
        return f"{size:.1f} {units[unit_index]}"

    @staticmethod
    def _format_duration(seconds: float) -> str:
        total_seconds = int(max(seconds, 0))
        minutes, secs = divmod(total_seconds, 60)
        hours, minutes = divmod(minutes, 60)
        if hours:
            return f"{hours}ч {minutes}м {secs}с"
        if minutes:
            return f"{minutes}м {secs}с"
        return f"{secs}с"

    def extract_document_text(self, file_path: Path) -> str:
        suffix = file_path.suffix.lower()

        if suffix in {".txt", ".md", ".markdown", ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs", ".c", ".cpp", ".h", ".hpp", ".css", ".scss", ".log", ".ini", ".cfg", ".sql", ".srt", ".rtf", ".yaml", ".yml"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")

        if suffix == ".json":
            raw = json.loads(file_path.read_text(encoding="utf-8", errors="ignore"))
            return json.dumps(raw, ensure_ascii=False, indent=2)

        if suffix == ".csv":
            rows: list[str] = []
            with file_path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
                reader = csv.reader(handle)
                for row in reader:
                    rows.append(" | ".join(cell.strip() for cell in row))
            return "\n".join(rows)

        if suffix in {".html", ".htm", ".xml"}:
            soup = BeautifulSoup(file_path.read_text(encoding="utf-8", errors="ignore"), "lxml")
            return soup.get_text("\n", strip=True)

        if suffix == ".pdf":
            reader = PdfReader(str(file_path))
            return "\n\n".join((page.extract_text() or "").strip() for page in reader.pages)

        if suffix == ".docx":
            doc = Document(str(file_path))
            lines = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    lines.append(" | ".join(cell.text.strip() for cell in row.cells))
            return "\n".join(lines)

        if suffix == ".pptx":
            presentation = Presentation(str(file_path))
            lines: list[str] = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                lines.append(f"Slide {slide_index}")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        lines.append(text.strip())
            return "\n".join(lines)

        if suffix == ".xlsx":
            workbook = load_workbook(filename=str(file_path), data_only=True)
            lines: list[str] = []
            for sheet in workbook.worksheets:
                lines.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                    if values:
                        lines.append(" | ".join(values))
            return "\n".join(lines)

        if suffix in self.TEXT_EXTENSIONS:
            return file_path.read_text(encoding="utf-8", errors="ignore")

        return ""

    def extract_video_keyframes(self, video_path: Path, output_dir: Path, max_frames: int) -> list[Path]:
        capture = cv2.VideoCapture(str(video_path))
        if not capture.isOpened():
            return []

        frame_count = int(capture.get(cv2.CAP_PROP_FRAME_COUNT))
        if frame_count <= 0:
            capture.release()
            return []

        sample_indexes = sorted({int(frame_count * index / max(max_frames, 1)) for index in range(max_frames)})
        sample_indexes = [min(index, frame_count - 1) for index in sample_indexes]

        result: list[Path] = []
        for position, frame_index in enumerate(sample_indexes, start=1):
            capture.set(cv2.CAP_PROP_POS_FRAMES, frame_index)
            ok, frame = capture.read()
            if not ok:
                continue
            output_path = output_dir / f"frame_{position:02d}.jpg"
            cv2.imwrite(str(output_path), frame)
            result.append(output_path)

        capture.release()
        return result

    def extract_audio_from_video(self, video_path: Path, output_dir: Path) -> Path | None:
        output_path = output_dir / "video_audio.mp3"
        command = [
            self.settings.ffmpeg_binary,
            "-loglevel",
            "error",
            "-y",
            "-i",
            str(video_path),
            "-vn",
            "-ac",
            "1",
            "-ar",
            "16000",
            str(output_path),
        ]
        try:
            subprocess.run(command, check=True, capture_output=True, text=True)
        except FileNotFoundError:
            LOGGER.exception("Анализ: ffmpeg не найден при извлечении аудио из видео %s", video_path.name)
            return None
        except subprocess.CalledProcessError:
            LOGGER.exception("Анализ: ffmpeg не смог извлечь аудио из видео %s", video_path.name)
            return None
        return output_path if output_path.exists() else None

    def split_audio_if_needed(self, audio_path: Path, output_dir: Path) -> list[Path]:
        segment_seconds = max(int(self.settings.audio_segment_seconds), 0)
        if segment_seconds <= 0:
            LOGGER.info("Транскрипция: сегментация отключена, использую исходный файл %s", audio_path.name)
            return [audio_path]

        duration_seconds = self._probe_media_duration_seconds(audio_path)
        if duration_seconds is not None and duration_seconds <= segment_seconds:
            LOGGER.info(
                "Транскрипция: файл %s короткий (%s), сегментация не нужна",
                audio_path.name,
                self._format_duration(duration_seconds),
            )
            return [audio_path]

        LOGGER.info(
            "Транскрипция: делю файл %s на сегменты по %s",
            audio_path.name,
            self._format_duration(segment_seconds),
        )

        segment_dir = output_dir / "segments"
        segment_dir.mkdir(parents=True, exist_ok=True)
        pattern = segment_dir / "segment_%03d.mp3"
        command = [
            self.settings.ffmpeg_binary,
            "-loglevel",
            "error",
            "-y",
            "-i",
            str(audio_path),
            "-f",
            "segment",
            "-segment_time",
            str(segment_seconds),
            "-ac",
            "1",
            "-ar",
            "16000",
            "-b:a",
            "64k",
            str(pattern),
        ]
        try:
            subprocess.run(command, check=True, capture_output=True, text=True)
        except FileNotFoundError as exc:
            raise RuntimeError("ffmpeg is required to split long audio files.") from exc
        except subprocess.CalledProcessError as exc:
            raise RuntimeError(exc.stderr.strip() or "ffmpeg failed to split the audio file.") from exc

        segments = sorted(segment_dir.glob("segment_*.mp3"))
        if not segments:
            raise RuntimeError("Failed to create audio segments.")
        LOGGER.info("Транскрипция: создано сегментов %s для файла %s", len(segments), audio_path.name)
        return segments

    def _probe_media_duration_seconds(self, media_path: Path) -> float | None:
        command = [
            self.settings.ffprobe_binary,
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "json",
            str(media_path),
        ]
        try:
            result = subprocess.run(command, check=True, capture_output=True, text=True)
        except (FileNotFoundError, subprocess.CalledProcessError):
            return None

        try:
            payload = json.loads(result.stdout or "{}")
            duration = float(payload.get("format", {}).get("duration"))
        except (TypeError, ValueError, json.JSONDecodeError):
            return None
        return duration if duration > 0 else None

    def _build_video_download_path(self, message: object) -> Path:
        chat_id = int(getattr(message, "chat_id"))
        message_id = int(getattr(message, "id"))
        message_date = getattr(message, "date", None)
        date_part = message_date.date().isoformat() if message_date else "unknown-date"

        file_name = self.get_file_name(message, "video") or f"video_{message_id}.mp4"
        safe_file_name = self._sanitize_file_name(file_name)
        if "." not in safe_file_name:
            extension = getattr(getattr(message, "file", None), "ext", None) or ".mp4"
            safe_file_name += extension

        target_dir = self.settings.video_download_dir / str(chat_id) / date_part
        target_dir.mkdir(parents=True, exist_ok=True)
        return target_dir / f"{message_id}_{safe_file_name}"

    @staticmethod
    def _sanitize_file_name(file_name: str) -> str:
        cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", file_name).strip(" .")
        return cleaned[:180] or "video.mp4"





